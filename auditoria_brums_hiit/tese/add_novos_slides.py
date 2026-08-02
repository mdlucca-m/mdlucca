# -*- coding: utf-8 -*-
"""Anexa à apresentação os slides das análises novas, no mesmo estilo do deck:
  • Invariância estrita (resíduos) e parcial — fecha a hierarquia de invariância.
  • Carga interna × perfil de humor e preditores de estado de fadiga.
Idempotente por título. Uso: python add_novos_slides.py  (de dentro de tese/)."""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE=os.path.dirname(os.path.abspath(__file__)); ROOT=os.path.dirname(HERE)
DECK=os.path.join(HERE,'Apresentacao_BRUMS_HIIT.pptx')
prs=Presentation(DECK)
SW,SH=prs.slide_width,prs.slide_height

TEAL='0E8C86'; INK='14243A'; CARDBG='F5F8FC'; CARDBD='DCE4EE'; GRAY='5B6B82'
GOLD='C7871B'; FAINT='93A1B5'; BLUE='245C8B'; CORAL='C0392B'; VIOLET='7A4FB0'; GREEN='2E7D5B'
hexc=lambda c:RGBColor.from_string(c)
LAYOUT=prs.slides[13].slide_layout

def has_title(t):
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and t in sh.text_frame.text: return True
    return False

def blank_slide():
    s=prs.slides.add_slide(LAYOUT)
    for ph in list(s.placeholders): ph._element.getparent().remove(ph._element)
    return s

def textbox(slide,x,y,cx,cy,anchor=MSO_ANCHOR.MIDDLE):
    tb=slide.shapes.add_textbox(Emu(x),Emu(y),Emu(cx),Emu(cy)); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0; tf.vertical_anchor=anchor; return tf

def run(p,txt,sz,color,bold=False,ital=False,font='Calibri',spc=None):
    r=p.add_run(); r.text=txt; f=r.font; f.size=Pt(sz); f.bold=bold; f.italic=ital; f.name=font; f.color.rgb=hexc(color)
    if spc is not None: r._r.get_or_add_rPr().set('spc',str(spc))

def card(slide,x,y,cx,cy,head,desc,col):
    rect=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Emu(x),Emu(y),Emu(cx),Emu(cy))
    rect.adjustments[0]=0.09; rect.fill.solid(); rect.fill.fore_color.rgb=hexc(CARDBG)
    rect.line.color.rgb=hexc(CARDBD); rect.line.width=Pt(1); rect.shadow.inherit=False
    tf=textbox(slide,x+160000,y+55000,cx-320000,cy-110000,MSO_ANCHOR.TOP)
    run(tf.paragraphs[0],head,13,col,bold=True,font='Cambria')
    run(tf.add_paragraph(),desc,10.5,GRAY)

def fig_slide(eyebrow,title,img,cards,takeaway):
    s=blank_slide()
    tf=textbox(s,502920,384048,10972800,292608); run(tf.paragraphs[0],eyebrow,12,TEAL,bold=True,spc=200)
    tf=textbox(s,457200,676656,11247120,860000); run(tf.paragraphs[0],title,26,INK,bold=True,font='Cambria')
    # figura à esquerda, dimensionada por altura máx preservando aspecto
    w,h=Image.open(img).size; ratio=h/w
    maxH=3760000; maxW=6850000
    fw=maxW; fh=int(fw*ratio)
    if fh>maxH: fh=maxH; fw=int(fh/ratio)
    fx=502920; fy=1560000+(maxH-fh)//2
    s.shapes.add_picture(img,Emu(fx),Emu(fy),Emu(fw),Emu(fh))
    # cards à direita
    x0=7560000; cx=SW-x0-502920; n=len(cards); gap=150000
    ch=int((3760000-(n-1)*gap)/n)
    for i,(head,desc,col) in enumerate(cards):
        card(s,x0,1560000+i*(ch+gap),cx,ch,head,desc,col)
    # barra de destaque
    ty=5560000; bw=SW-2*548640
    rect=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Emu(548640),Emu(ty),Emu(bw),Emu(700000))
    rect.adjustments[0]=0.12; rect.fill.solid(); rect.fill.fore_color.rgb=hexc(CARDBG)
    rect.line.color.rgb=hexc(CARDBD); rect.line.width=Pt(1); rect.shadow.inherit=False
    tf=textbox(s,777240,ty,bw-450000,700000); p=tf.paragraphs[0]
    run(p,'▸  ',13,GOLD,bold=True); run(p,takeaway,13,INK,bold=True)
    # rodapé
    tf=textbox(s,457200,6473952,9000000,274320); run(tf.paragraphs[0],'BRUMS × HIIT · síntese analítica',9,FAINT)
    tf=textbox(s,11247120,6473952,548640,274320); tf.paragraphs[0].alignment=PP_ALIGN.RIGHT
    run(tf.paragraphs[0],str(len(prs.slides._sldIdLst)),9,FAINT)
    return s

added=[]
T1='Invariância estrita (resíduos) e parcial'
if not has_title(T1):
    fig_slide('RESULTADOS · INVARIÂNCIA DE MEDIDA', T1,
        os.path.join(ROOT,'invariancia_multigrupo','invariancia_estrita_parcial_fig.png'),
        [('Métrica sustentada','ΔCFI −0,004 (modelo conjunto); Tucker φ 0,992 — as cargas são congruentes pré→pós.',BLUE),
         ('Escalar aproximada','viés residual de intercepto RMS 0,135; deslocamento latente na fadiga (+0,56) e vigor (−0,26).',GREEN),
         ('Estrita no limite','fixando também as variâncias de erro θ: CFI 0,913→0,898 (ΔCFI −0,015).',GOLD),
         ('Parcial restabelece','liberando “Sonolento” (fadiga_3) e vigor_4, o viés cai 0,135→0,064 — quebra local, não global.',CORAL)],
        'A mudança pré→pós é um deslocamento de estado, não artefato de medida: a invariância se sustenta e a quebra é local e identificada.')
    added.append(T1)

T2='Carga interna × humor e preditores de estado de fadiga'
if not has_title(T2):
    fig_slide('RESULTADOS · CARGA × HUMOR E PREDITORES', T2,
        os.path.join(ROOT,'carga_humor','carga_humor_fig.png'),
        [('Desacoplamento','PSE/FC/TRIMP × humor: nada sobrevive à FDR — a carga prescreve o estímulo, não a resposta.',BLUE),
         ('Mais sensível','fadiga física: AUC 0,90 (estado-lábil, ICC 0,28) — a melhor sentinela de estado.',CORAL),
         ('Sensível E confiável','fadiga mental (ICC 0,72) e depressão (ICC 0,79) — marcadores estáveis dia a dia.',GREEN),
         ('Confiável, porém cega','tensão: a mais estável (ICC 0,83), mas AUC 0,54 — não discrimina fadiga.',GOLD)],
        'Para monitorar fadiga: fadiga física pela sensibilidade; fadiga mental e depressão pela confiabilidade. Reencontra-se o eixo energia–fadiga.')
    added.append(T2)

T3='Estatística descritiva: distribuições e normalidade'
if not has_title(T3):
    fig_slide('RESULTADOS · ESTATÍSTICA DESCRITIVA', T3,
        os.path.join(ROOT,'descritivas_testes','descritiva_completa_fig.png'),
        [('Histogramas','distribuição da fadiga física pré vs. pós — deslocamento à direita no pós.',BLUE),
         ('Box plots','pré×pós das variáveis-chave: fadiga sobe, vigor recua (eixo energia–fadiga).',CORAL),
         ('Dispersão','vigor × fadiga com r≈−0,44 — a relação inversa que define o eixo.',GREEN),
         ('Normalidade (Q–Q)','desvios da normal no pré e no pós (Shapiro–Wilk) — justifica permutação.',GOLD)],
        'Distribuições assimétricas e não-normais (efeito piso nas negativas): as decisões são confirmadas por testes não-paramétricos e permutação.')
    added.append(T3)

T4='Protocolo de HIIT: 4 × 4 min a 104% da velocidade de pico'
if not has_title(T4):
    fig_slide('MÉTODO · PROTOCOLO DE HIIT', T4,
        os.path.join(ROOT,'dias_hiit','hiit_protocolo_fig.png'),
        [('Estímulo','4 séries de 4 min a 104% da velocidade de pico (teste de campo), 3 min de intervalo.',BLUE),
         ('FC das séries','pico ~186 bpm; média das 4 séries ~179 bpm — próximo ao teto aeróbio.',CORAL),
         ('Recuperação','FC de recuperação sobe ao longo da sessão (~124 bpm): deriva/recuperação incompleta.',GREEN),
         ('PSE','esforço percebido final ~9,9 (quase o teto da escala 0–10).',GOLD)],
        'Carga interna quase-máxima e progressiva — o regime de teto explica por que a carga não prediz a resposta de humor.')
    added.append(T4)

if added:
    prs.save(DECK)
    print('slides adicionados:',added,'| total:',len(prs.slides._sldIdLst))
else:
    print('slides já presentes; nada a fazer. total:',len(prs.slides._sldIdLst))

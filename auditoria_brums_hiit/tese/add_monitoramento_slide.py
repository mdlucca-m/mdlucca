# -*- coding: utf-8 -*-
"""Anexa à apresentação um slide de síntese visual com o painel de
monitoramento (gauge · radar · monitoramento diário · bolhas 4D), no mesmo
estilo do deck. Idempotente: se o slide já existir, não faz nada.
Uso: python add_monitoramento_slide.py  (de dentro de tese/)
Requer a figura tese/figuras/monitoramento_viz.png (ver monitoramento_viz.py)."""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
DECK = os.path.join(HERE, 'Apresentacao_BRUMS_HIIT.pptx')
IMG = os.path.join(HERE, 'figuras', 'monitoramento_viz.png')
TITLE = 'Painel de monitoramento e visualizações'

prs = Presentation(DECK)

# idempotência: já existe o slide?
for s in prs.slides:
    for sh in s.shapes:
        if sh.has_text_frame and TITLE in sh.text_frame.text:
            print('slide já presente; nada a fazer.')
            raise SystemExit(0)

layout = prs.slides[13].slide_layout  # mesmo layout dos slides de figura
slide = prs.slides.add_slide(layout)
for ph in list(slide.placeholders):
    ph._element.getparent().remove(ph._element)

TEAL='0E8C86'; INK='14243A'; CARDBG='F5F8FC'; CARDBD='DCE4EE'; GRAY='5B6B82'
GOLD='C7871B'; FAINT='93A1B5'; BLUE='245C8B'; CORAL='C0392B'; VIOLET='7A4FB0'
hexc = lambda c: RGBColor.from_string(c)

def textbox(x, y, cx, cy, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(cx), Emu(cy))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf

def run(p, txt, sz, color, bold=False, ital=False, font='Calibri', spc=None):
    r = p.add_run(); r.text = txt; f = r.font
    f.size = Pt(sz); f.bold = bold; f.italic = ital; f.name = font
    f.color.rgb = hexc(color)
    if spc is not None:
        r._r.get_or_add_rPr().set('spc', str(spc))

tf = textbox(502920, 384048, 10972800, 292608)
run(tf.paragraphs[0], 'RESULTADOS · SÍNTESE VISUAL', 12, TEAL, bold=True, spc=200)
tf = textbox(457200, 676656, 11247120, 914400)
run(tf.paragraphs[0], TITLE, 29, INK, bold=True, font='Cambria')

slide.shapes.add_picture(IMG, Emu(502920), Emu(1600200), Emu(5394960), Emu(4105000))

cards = [
 ('A · Medidores (gauge)', 'indicadores-chave sobre zonas de referência: dz fadiga física, AUC, perfil iceberg e HTMT', BLUE),
 ('B · Radar do humor', 'perfil das 6 subescalas pré vs. pós — a “montanha” de vigor achata (erosão do iceberg)', CORAL),
 ('C · Monitoramento diário', 'carga de humor D1→D7 com faixas de alerta e dias de HIIT — pico no D7', GOLD),
 ('D · Mapa 4D das variáveis', 'resposta aguda × acúmulo × individualidade × consenso direcional por variável', VIOLET),
]
cx = 5532120; cardh = 900000; gap = 130000; x0 = 6126480; y0 = 1600200
for i, (head, desc, col) in enumerate(cards):
    y = y0 + i * (cardh + gap)
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(x0), Emu(y), Emu(cx), Emu(cardh))
    rect.adjustments[0] = 0.09
    rect.fill.solid(); rect.fill.fore_color.rgb = hexc(CARDBG)
    rect.line.color.rgb = hexc(CARDBD); rect.line.width = Pt(1); rect.shadow.inherit = False
    tf = textbox(x0 + 228600, y + 60000, cx - 457200, cardh - 120000)
    run(tf.paragraphs[0], head, 15, col, bold=True, font='Cambria')
    run(tf.add_paragraph(), desc, 11.5, GRAY)

ty = 5806440
rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(548640), Emu(ty), Emu(11109960), Emu(658368))
rect.adjustments[0] = 0.125
rect.fill.solid(); rect.fill.fore_color.rgb = hexc(CARDBG)
rect.line.color.rgb = hexc(CARDBD); rect.line.width = Pt(1); rect.shadow.inherit = False
tf = textbox(777240, ty, 10698480, 658368)
p = tf.paragraphs[0]
run(p, '▸  ', 13.5, GOLD, bold=True)
run(p, 'Um só painel resume as três âncoras do estudo: resposta no eixo energia–fadiga, forte '
       'individualidade entre atletas e monitoramento por tendência (a fadiga física é a sentinela).',
    13.5, INK, bold=True)

tf = textbox(457200, 6473952, 8229600, 274320)
run(tf.paragraphs[0], 'Síntese visual · monitoramento', 9, FAINT)
tf = textbox(11247120, 6473952, 548640, 274320)
tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
run(tf.paragraphs[0], str(len(prs.slides._sldIdLst)), 9, FAINT)

prs.save(DECK)
print('slide adicionado. total de slides:', len(prs.slides._sldIdLst))
